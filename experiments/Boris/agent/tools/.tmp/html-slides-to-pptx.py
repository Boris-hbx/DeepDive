#!/usr/bin/env python3
"""
html-slides-to-pptx — Convert an HTML slide deck to a PPTX with full-page screenshots.

Splits each slide into a standalone HTML, screenshots via Edge headless, embeds into PPTX.

Usage:
    python html-slides-to-pptx.py <input.html> [output.pptx]
"""

import json
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu

EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
if not Path(EDGE).exists():
    EDGE = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
WIDTH_PX = 1920
HEIGHT_PX = 1080


def extract_slides(html_path: Path) -> tuple[str, str, list[str]]:
    """Extract head/style and individual slide divs from the HTML."""
    raw = html_path.read_text(encoding="utf-8")

    head_m = re.search(r"(<head.*?</head>)", raw, re.DOTALL)
    head = head_m.group(1) if head_m else "<head></head>"

    canvas_line = '<canvas id="particles" style="position:fixed;top:0;left:0;width:100%;height:100%;pointer-events:none;z-index:0;opacity:0.4;"></canvas>'

    slides = re.findall(
        r'(<div\s+class="slide[^"]*"[^>]*data-slide="\d+"[^>]*>.*?</div>\s*\n)',
        raw, re.DOTALL,
    )
    if not slides:
        slides = re.findall(
            r'(<div\s+class="slide[^"]*"[^>]*>.*?</div>\s*(?=<div\s+class="slide|</div>\s*<div\s+class="nav))',
            raw, re.DOTALL,
        )

    script_m = re.search(r"(<script>.*?// Particle background.*?animateParticles\(\);)", raw, re.DOTALL)
    particle_script = script_m.group(1) if script_m else ""

    return head, canvas_line, particle_script, slides


def build_single_slide_html(head: str, canvas_line: str, particle_script: str, slide_div: str) -> str:
    """Build a standalone HTML for one slide, forced to active/visible."""
    slide_div = slide_div.replace('class="slide"', 'class="slide active"', 1)
    if 'class="slide active"' not in slide_div:
        slide_div = re.sub(r'class="slide([^"]*)"', r'class="slide active\1"', slide_div, count=1)

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
{head}
<body>
{canvas_line}
<div class="deck" style="width:100vw;height:100vh;position:relative;z-index:1;">
{slide_div}
</div>
{particle_script}
</body>
</html>"""


def screenshot(html_content: str, output_png: Path):
    """Write HTML to temp file, screenshot with Edge headless."""
    tmp_html = output_png.with_suffix(".html")
    tmp_html.write_text(html_content, encoding="utf-8")

    cmd = [
        EDGE,
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        f"--screenshot={output_png}",
        f"--window-size={WIDTH_PX},{HEIGHT_PX}",
        f"file:///{tmp_html.as_posix()}",
    ]
    subprocess.run(cmd, capture_output=True, timeout=20)
    tmp_html.unlink(missing_ok=True)

    return output_png.exists()


def build_pptx(png_files: list[Path], output_path: Path):
    """Create a PPTX with each PNG as a full-slide image."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for png in png_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(
            str(png),
            Emu(0), Emu(0),
            SLIDE_W, SLIDE_H,
        )

    prs.save(str(output_path))


def main():
    if len(sys.argv) < 2:
        print("Usage: python html-slides-to-pptx.py <input.html> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    html_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else html_path.with_suffix(".pptx")

    print(f"Reading {html_path}...")
    head, canvas, particle_js, slides = extract_slides(html_path)
    print(f"Found {len(slides)} slides")

    if not slides:
        print("ERROR: No slides found in HTML", file=sys.stderr)
        sys.exit(1)

    tmp_dir = Path(tempfile.mkdtemp(prefix="slides_"))
    png_files = []

    for i, slide_div in enumerate(slides):
        png_path = tmp_dir / f"slide_{i:02d}.png"
        html = build_single_slide_html(head, canvas, particle_js, slide_div)
        print(f"  Screenshotting slide {i+1}/{len(slides)}...", end=" ", flush=True)
        if screenshot(html, png_path):
            png_files.append(png_path)
            size_kb = png_path.stat().st_size // 1024
            print(f"OK ({size_kb}KB)")
        else:
            print("FAILED")

    print(f"\nBuilding PPTX with {len(png_files)} slides...")
    build_pptx(png_files, output_path)
    print(f"Saved: {output_path}")

    for f in tmp_dir.glob("*"):
        f.unlink()
    tmp_dir.rmdir()

    print(json.dumps({"status": "ok", "path": str(output_path), "slides": len(png_files)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
