#!/usr/bin/env python3
"""
html-slides-to-pptx — Screenshot each slide from an HTML deck via CDP, assemble into PPTX.

Uses Edge headless + Chrome DevTools Protocol over WebSocket to:
1. Load the original HTML with all CSS/JS/SVG/animations intact
2. Call show(i) to navigate to each slide
3. Capture high-DPI screenshots (2x)
4. Embed as full-page images in PPTX

Usage:
    python html-slides-to-pptx.py <input.html> [output.pptx]
"""

import asyncio
import base64
import json
import shutil
import socket
import subprocess
import sys
import tempfile
import time
import urllib.request
from pathlib import Path

import websockets

from pptx import Presentation
from pptx.util import Inches, Emu

EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
if not Path(EDGE).exists():
    EDGE = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
WIDTH_PX = 1920
HEIGHT_PX = 1080
DPR = 2


def find_free_port():
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("", 0))
        return s.getsockname()[1]


async def capture_slides(html_path: Path, out_dir: Path) -> list[Path]:
    port = find_free_port()
    user_data = Path(tempfile.mkdtemp(prefix="edge_cdp_"))

    cmd = [
        EDGE,
        f"--remote-debugging-port={port}",
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--disable-extensions",
        "--no-first-run",
        "--no-default-browser-check",
        f"--user-data-dir={user_data}",
        f"--window-size={WIDTH_PX},{HEIGHT_PX}",
        f"file:///{html_path.resolve().as_posix()}",
    ]

    proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

    # Wait for CDP to be ready
    ws_url = None
    for _ in range(15):
        await asyncio.sleep(1)
        try:
            with urllib.request.urlopen(f"http://127.0.0.1:{port}/json", timeout=2) as resp:
                targets = json.loads(resp.read())
            for t in targets:
                if t.get("type") == "page":
                    ws_url = t["webSocketDebuggerUrl"]
                    break
            if ws_url:
                break
        except Exception:
            continue

    if not ws_url:
        proc.kill()
        raise RuntimeError("Could not connect to Edge CDP")

    msg_id = 0

    async def send_cdp(ws, method, params=None):
        nonlocal msg_id
        msg_id += 1
        payload = {"id": msg_id, "method": method}
        if params:
            payload["params"] = params
        await ws.send(json.dumps(payload))
        while True:
            resp = json.loads(await ws.recv())
            if resp.get("id") == msg_id:
                return resp.get("result", {})

    async with websockets.connect(ws_url, max_size=50 * 1024 * 1024) as ws:
        # Wait for page to fully load
        await asyncio.sleep(2)

        # Set high-DPI viewport
        await send_cdp(ws, "Emulation.setDeviceMetricsOverride", {
            "width": WIDTH_PX,
            "height": HEIGHT_PX,
            "deviceScaleFactor": DPR,
            "mobile": False,
        })
        await asyncio.sleep(0.5)

        # Get slide count
        r = await send_cdp(ws, "Runtime.evaluate", {
            "expression": "document.querySelectorAll('.slide').length",
            "returnByValue": True,
        })
        total = r.get("result", {}).get("value", 0)
        print(f"  Found {total} slides")

        png_files = []
        for i in range(total):
            # Navigate to slide
            await send_cdp(ws, "Runtime.evaluate", {
                "expression": f"show({i})",
            })
            await asyncio.sleep(0.6)

            # Capture screenshot
            r = await send_cdp(ws, "Page.captureScreenshot", {
                "format": "png",
                "captureBeyondViewport": False,
            })

            png_data = base64.b64decode(r.get("data", ""))
            png_path = out_dir / f"slide_{i:02d}.png"
            png_path.write_bytes(png_data)
            png_files.append(png_path)
            print(f"  Slide {i+1}/{total}: {len(png_data)//1024}KB")

    proc.kill()
    proc.wait()
    shutil.rmtree(user_data, ignore_errors=True)
    return png_files


def build_pptx(png_files: list[Path], output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for png in png_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png), Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    prs.save(str(output_path))


def main():
    if len(sys.argv) < 2:
        print("Usage: python html2pptx.py <input.html> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    html_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else html_path.with_suffix(".pptx")

    print(f"Converting {html_path.name} -> {output_path.name}")

    tmp_dir = Path(tempfile.mkdtemp(prefix="slides_"))
    png_files = asyncio.run(capture_slides(html_path, tmp_dir))

    print(f"\nBuilding PPTX with {len(png_files)} slides...")
    build_pptx(png_files, output_path)
    print(f"Done: {output_path} ({output_path.stat().st_size // 1024}KB)")

    for f in tmp_dir.glob("*"):
        f.unlink()
    tmp_dir.rmdir()


if __name__ == "__main__":
    main()
