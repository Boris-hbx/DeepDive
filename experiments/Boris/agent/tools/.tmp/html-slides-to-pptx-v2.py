#!/usr/bin/env python3
"""
html-slides-to-pptx — Screenshot each slide from an HTML deck and assemble into PPTX.

Uses Edge DevTools Protocol (CDP) to navigate slides in the original HTML,
preserving all animations, SVG diagrams, fonts, and visual effects.

Usage:
    python html-slides-to-pptx.py <input.html> [output.pptx]
"""

import json
import os
import subprocess
import sys
import tempfile
import time
import base64
import socket
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu

EDGE = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
if not Path(EDGE).exists():
    EDGE = r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
WIDTH_PX = 1920
HEIGHT_PX = 1080
DPR = 2  # device pixel ratio for crisp screenshots


def find_free_port():
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("", 0))
        return s.getsockname()[1]


def cdp_send(ws_url, method, params=None):
    """Minimal CDP over HTTP — use /json/protocol endpoint instead of WebSocket."""
    pass


def screenshot_all_slides(html_path: Path, out_dir: Path) -> list[Path]:
    """Launch Edge, load HTML, iterate slides via JS, screenshot each via CDP."""
    port = find_free_port()
    user_data = Path(tempfile.mkdtemp(prefix="edge_cdp_"))

    cmd = [
        EDGE,
        f"--remote-debugging-port={port}",
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--disable-extensions",
        "--no-first-run",
        "--no-default-browser-check",
        f"--user-data-dir={user_data}",
        f"--window-size={WIDTH_PX},{HEIGHT_PX}",
        f"file:///{html_path.resolve().as_posix()}",
    ]

    proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    time.sleep(3)  # wait for Edge to start

    # Get CDP WebSocket URL
    targets_url = f"http://127.0.0.1:{port}/json"
    for attempt in range(10):
        try:
            with urllib.request.urlopen(targets_url, timeout=3) as resp:
                targets = json.loads(resp.read())
            break
        except Exception:
            time.sleep(1)
    else:
        proc.kill()
        raise RuntimeError("Could not connect to Edge CDP")

    page_ws = None
    for t in targets:
        if t.get("type") == "page":
            page_ws = t.get("webSocketDebuggerUrl")
            page_id = t.get("id")
            break

    if not page_ws:
        proc.kill()
        raise RuntimeError("No page target found")

    # Use websocket via a simple Python implementation
    import websocket_cdp
    ws = websocket_cdp.CDPConnection(page_ws)

    # Wait for page load
    time.sleep(2)

    # Set device metrics for high-DPI
    ws.send_command("Emulation.setDeviceMetricsOverride", {
        "width": WIDTH_PX,
        "height": HEIGHT_PX,
        "deviceScaleFactor": DPR,
        "mobile": False,
    })
    time.sleep(0.5)

    # Get total slide count
    result = ws.send_command("Runtime.evaluate", {
        "expression": "document.querySelectorAll('.slide').length"
    })
    total = result.get("result", {}).get("value", 0)
    print(f"  Found {total} slides in browser")

    png_files = []
    for i in range(total):
        # Navigate to slide i
        ws.send_command("Runtime.evaluate", {
            "expression": f"show({i})"
        })
        time.sleep(0.8)  # wait for transition animation

        # Capture screenshot
        result = ws.send_command("Page.captureScreenshot", {
            "format": "png",
            "quality": 100,
            "clip": {
                "x": 0, "y": 0,
                "width": WIDTH_PX, "height": HEIGHT_PX,
                "scale": 1,
            },
            "captureBeyondViewport": False,
        })

        png_data = base64.b64decode(result.get("data", ""))
        png_path = out_dir / f"slide_{i:02d}.png"
        png_path.write_bytes(png_data)
        size_kb = len(png_data) // 1024
        png_files.append(png_path)
        print(f"  Slide {i+1}/{total}: {size_kb}KB")

    ws.close()
    proc.kill()

    # Cleanup user data dir
    import shutil
    try:
        shutil.rmtree(user_data, ignore_errors=True)
    except Exception:
        pass

    return png_files


def build_pptx(png_files: list[Path], output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for png in png_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png), Emu(0), Emu(0), SLIDE_W, SLIDE_H)

    prs.save(str(output_path))


def main():
    if len(sys.argv) < 2:
        print("Usage: python html-slides-to-pptx.py <input.html> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    html_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else html_path.with_suffix(".pptx")

    print(f"Converting {html_path.name} → {output_path.name}")

    tmp_dir = Path(tempfile.mkdtemp(prefix="slides_"))

    png_files = screenshot_all_slides(html_path, tmp_dir)
    print(f"\nBuilding PPTX with {len(png_files)} slides...")
    build_pptx(png_files, output_path)
    print(f"Saved: {output_path} ({output_path.stat().st_size // 1024}KB)")

    # Cleanup
    for f in tmp_dir.glob("*"):
        f.unlink()
    tmp_dir.rmdir()


if __name__ == "__main__":
    main()
