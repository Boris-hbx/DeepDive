#!/usr/bin/env python3
"""
read-pptx — Extract content from a PowerPoint file as structured JSON.

Usage:
    python read-pptx.py <input.pptx>

Output JSON format:
    {
      "slide_count": 5,
      "slides": [
        {
          "number": 1,
          "title": "Slide Title",
          "texts": ["text block 1", "text block 2"],
          "notes": "Speaker notes if any",
          "shapes_count": 3,
          "images_count": 1
        }
      ]
    }
"""

import json
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_pptx(filepath):
    prs = Presentation(filepath)
    result = {
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides, 1):
        slide_data = {
            "number": i,
            "title": "",
            "texts": [],
            "notes": "",
            "shapes_count": len(slide.shapes),
            "images_count": 0,
        }

        # Title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text

        # All text content
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data["images_count"] += 1
                continue

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != slide_data["title"]:
                    slide_data["texts"].append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    slide_data["texts"].append({"table": rows})

        # Notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        result["slides"].append(slide_data)

    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python read-pptx.py <input.pptx>", file=sys.stderr)
        sys.exit(1)

    filepath = sys.argv[1]
    if not Path(filepath).exists():
        print(f"Error: file not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    try:
        result = read_pptx(filepath)
    except Exception as e:
        print(f"Error reading PPTX: {e}", file=sys.stderr)
        sys.exit(1)

    print(json.dumps(result, ensure_ascii=False, indent=2))
