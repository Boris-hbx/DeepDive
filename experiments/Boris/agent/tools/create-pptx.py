#!/usr/bin/env python3
"""
create-pptx — Generate a styled PowerPoint file from JSON content.

Style: 微软雅黑, red title with underline, high density, max 18pt body.

Usage:
    python create-pptx.py <json_file_or_stdin> [output_path]

Input JSON format:
    {
      "title": "Presentation Title",
      "subtitle": "Optional subtitle",
      "slides": [
        {
          "title": "Slide Title",
          "content": "Body text",
          "bullets": ["point 1", "point 2"],
          "notes": "Speaker notes (optional)",
          "layout": "title|content|section|blank"
        }
      ]
    }
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Style constants ──────────────────────────────────────────────────────────
FONT_NAME = "微软雅黑"
FONT_NAME_EN = "Microsoft YaHei"
TITLE_SIZE = Pt(24)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(16)
BULLET_SIZE = Pt(15)
NOTES_SIZE = Pt(11)

COLOR_TITLE = RGBColor(0xC0, 0x00, 0x00)  # dark red
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_RED_LINE = RGBColor(0xC0, 0x00, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG_TITLE = RGBColor(0x1A, 0x1A, 0x2E)  # dark blue for title slide

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Cm(1.5)
MARGIN_T = Cm(1.2)
CONTENT_W = SLIDE_W - Cm(3.0)


def _set_font(run, size=BODY_SIZE, color=COLOR_BODY, bold=False):
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    # Set East Asian font
    rpr = run._r.get_or_add_rPr()
    ea = rpr.makeelement(qn('a:ea'), {})
    ea.set('typeface', FONT_NAME)
    rpr.append(ea)


def _add_red_line(slide, left, top, width):
    """Add a red horizontal line below the title."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, Pt(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_RED_LINE
    shape.line.fill.background()


def _add_title_box(slide, title_text, top=MARGIN_T):
    """Add a styled title text box with red color and underline."""
    txBox = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, Cm(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_font(run, size=TITLE_SIZE, color=COLOR_TITLE, bold=True)

    # Red line below title
    line_top = top + Cm(1.6)
    _add_red_line(slide, MARGIN_L, line_top, CONTENT_W)

    return line_top + Cm(0.6)


def _add_body_text(slide, text, top):
    """Add body text paragraph."""
    txBox = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, SLIDE_H - top - Cm(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=BODY_SIZE, color=COLOR_BODY)
    p.space_after = Pt(8)
    return txBox


def _add_bullets(slide, bullets, top):
    """Add bullet points with high density spacing."""
    txBox = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, SLIDE_H - top - Cm(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = f"▸ {bullet}"
        _set_font(run, size=BULLET_SIZE, color=COLOR_BODY)

    return txBox


def create_pptx(data, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_text = data.get("title", "Untitled")
    subtitle_text = data.get("subtitle", "")

    # ── Title slide (dark background) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG_TITLE

    # Title centered
    txBox = slide.shapes.add_textbox(Cm(2), Cm(2.5), SLIDE_W - Cm(4), Cm(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_font(run, size=Pt(36), color=COLOR_WHITE, bold=True)

    # Red line
    _add_red_line(slide, Cm(2), Cm(5.2), Cm(8))

    # Subtitle
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Cm(2), Cm(5.8), SLIDE_W - Cm(4), Cm(2))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        _set_font(run2, size=SUBTITLE_SIZE, color=RGBColor(0xCC, 0xCC, 0xCC))

    # ── Content slides ──
    for s in data.get("slides", []):
        layout_name = s.get("layout", "content")
        slide_title = s.get("title", "")

        if layout_name == "section":
            # Section divider — dark bg, large title
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = COLOR_BG_TITLE

            txBox = slide.shapes.add_textbox(Cm(2), Cm(2.8), SLIDE_W - Cm(4), Cm(2.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = slide_title
            _set_font(run, size=Pt(28), color=COLOR_WHITE, bold=True)

            _add_red_line(slide, Cm(2), Cm(5.0), Cm(6))

            if s.get("content"):
                txBox2 = slide.shapes.add_textbox(Cm(2), Cm(5.6), SLIDE_W - Cm(4), Cm(2))
                tf2 = txBox2.text_frame
                p2 = tf2.paragraphs[0]
                run2 = p2.add_run()
                run2.text = s["content"]
                _set_font(run2, size=BODY_SIZE, color=RGBColor(0xAA, 0xAA, 0xAA))

        elif layout_name == "blank":
            prs.slides.add_slide(prs.slide_layouts[6])

        else:
            # Standard content slide — white bg
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank for full control

            # Title + red line
            body_top = _add_title_box(slide, slide_title)

            # Content
            bullets = s.get("bullets", [])
            if bullets:
                _add_bullets(slide, bullets, body_top)
            elif s.get("content"):
                _add_body_text(slide, s["content"], body_top)

        # Speaker notes
        if s.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = s["notes"]

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create-pptx.py <json_file_or_-> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    src = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    if src == "-":
        data = json.loads(sys.stdin.read())
    else:
        with open(src, "r", encoding="utf-8") as f:
            data = json.load(f)

    result = create_pptx(data, output)
    print(json.dumps({"status": "ok", "path": str(result)}, ensure_ascii=False))
